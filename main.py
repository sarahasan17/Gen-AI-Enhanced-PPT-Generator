import os
from dotenv import load_dotenv
import pptx
import pandas as pd
import PyPDF2
import docx
import pdfplumber
import matplotlib.pyplot as plt
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import tempfile
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage
import re

# Load environment variables
load_dotenv()

# Retrieve API key from environment variable
def initialize_ai():
    return ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        groq_api_key=os.getenv("GROQ_API_KEY")
    )

llm = initialize_ai()

# Define custom formatting options
TITLE_FONT_SIZE = Pt(34)
SLIDE_FONT_SIZE = Pt(20)

# Function to generate title using AI
def generate_title_with_ai(text):
    prompt = f"""
    Generate a concise and professional PowerPoint presentation title based on the following content:
    {text}
    - The title should accurately reflect the main topic.
    - Keep it under 10 words.
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    return response.content.strip()

# Function to extract text and tables from PDFs
def extract_text_from_pdf(pdf_path):
    text = ""
    tables = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
            table = page.extract_table()
            if table:
                df = pd.DataFrame(table[1:], columns=table[0])
                tables.append(df)
    return text.strip(), tables

# Function to extract text and tables from Word documents
def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    tables = []
    doc = docx.Document(docx_path)
    text = []
    # Extract paragraph text
    for p in doc.paragraphs:
        text.append(p.text)
    for table in doc.tables:
        df = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append(pd.DataFrame(df))
    return text, tables

# Function to extract text and tables from PowerPoint slides
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    tables = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
            elif hasattr(shape, "table"):
                table = shape.table
                df = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                tables.append(pd.DataFrame(df))
    
    return text.strip(), tables

# Function to extract data from Excel
def extract_data_from_excel(excel_path):
    return pd.read_excel(excel_path)

def remove_unwanted_slides(prs):
    """Remove slides with generic notes or redundant content."""
    for slide in list(prs.slides):
        if slide.shapes.title and "Note:" in slide.shapes.title.text:
            text_content = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            if "The content of each slide should be concise and to the point" in text_content:
                prs.slides._sldIdLst.remove(slide._element)

# Function to enhance text using AI
def enhance_text_with_ai(text):
    prompt = f"""
    Improve the following content for a structured PowerPoint presentation:
    {text}
    - Generate the main content only directly without any AI responses.
    - Organize it into clear, professional sections with meaningful slide titles.
    - Summarize key points for a high-level overview.
    - Ensure slides contain concise bullet points, avoiding unnecessary symbols or long sentences.
    - Maintain a professional tone and formatting.
    - **Remove redundant or irrelevant** information from the text.
    - Present the data properly but do not write words like "key points" or "bullet points"
    - Avoid generic AI-generated responses like 'I'm ready to help' or 'Please provide content' or .
    - Avoid repeating data from Excel unless it adds value.
    - Avoid unecessary slides while presenting excel data in the beginning without effective content.
    - If there is no content worth generating do not geenrate any redundant data just ignore the slide.
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    return response.content.strip()

# Function to create slides with custom formatting
def add_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = str(title).strip()
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    
    for point in content.split("\n"):
        clean_point = point.strip().lstrip("• *+.-")
        if clean_point:
            p = text_frame.add_paragraph()
            p.text = f"{clean_point}"
            p.space_after = Inches(0.2)
            p.font.size = SLIDE_FONT_SIZE

def create_chart(prs, title, categories, values, chart_type):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = str(title).strip()
    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(str(title), values)  # Ensure title is a string
    
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
    slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)

def process_excel_data(prs, excel_data):
    if excel_data is not None and isinstance(excel_data, pd.DataFrame):
        # Drop duplicate rows and columns with mostly empty values
        excel_data = excel_data.drop_duplicates().dropna(axis=1, how='all')

        # Remove columns with only one unique value (not useful for presentation)
        excel_data = excel_data.loc[:, excel_data.nunique() > 1]

        # Skip slide creation if the data is empty after cleaning
        if excel_data.empty:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Excel Data Overview"
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()

        for index, row in excel_data.iterrows():
            p = text_frame.add_paragraph()
            p.text = " | ".join(map(str, row.values))
            p.font.size = SLIDE_FONT_SIZE

        # Generate charts for numerical columns
        numeric_cols = excel_data.select_dtypes(include=['number'])
        if not numeric_cols.empty:
            categories = excel_data.iloc[:, 0].astype(str).tolist()
            for col in numeric_cols.columns:
                values = excel_data[col].fillna(0).tolist()
                create_chart(prs, f"{col} Analysis", categories, values, XL_CHART_TYPE.COLUMN_CLUSTERED)

# Function to analyze text for graphs
def analyze_text_for_graphs(text, prs):
    prompt = f"""
    Analyze the following text and extract numerical data for charts.
    Output in JSON format: {text}
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    
    try:
        graph_data = eval(response.content)
        if "categories" in graph_data and "values" in graph_data:
            create_chart(prs, graph_data.get("title", "Generated Chart"), graph_data["categories"], graph_data["values"], chart_type=XL_CHART_TYPE.PIE)
    except:
        pass  # Skip if no valid graph data is found

def extract_trends_with_ai(text):
    prompt = f"""
    Identify trends in the following content and structure them into a table format:
    {text}
    - The table should contain columns for 'Year', 'Metric', and 'Value'.
    - Focus on numerical trends such as increases, decreases, and fluctuations over time.
    - Ensure the output follows the format: Year|Metric|Value\n2020|Revenue|5.0\n2023|Revenue|8.2.
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    return response.content.strip()

def create_trend_chart(prs, trend_data):
    try:
        rows = [x.split('|') for x in trend_data.strip().split('\n') if '|' in x]
        if not rows or len(rows[0]) != 3:
            return
        df = pd.DataFrame(rows, columns=['Year', 'Metric', 'Value'])
        df['Value'] = pd.to_numeric(df['Value'], errors='coerce')
        metrics = df['Metric'].unique()
        
        for metric in metrics:
            sub_df = df[df['Metric'] == metric]
            categories = sub_df['Year'].tolist()
            values = sub_df['Value'].tolist()
            create_chart(prs, metric, categories, values, XL_CHART_TYPE.LINE)
    except Exception as e:
        print(f"Error processing trend data: {e}")

# Function to create PowerPoint presentation
def create_ppt_from_text(text, tables=None, excel_data=None):
    enhanced_text = enhance_text_with_ai(text).strip()
    trend_data = extract_trends_with_ai(text)
    prs = Presentation()
    
    title = generate_title_with_ai(enhanced_text)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = str(title).strip()
    title_slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

    # Define phrases to filter out if no other content is present
    redundant_texts = [
        "The content of each slide should be concise",
        "redundant or irrelevant information should be avoided"
    ]

    sections = [s.strip() for s in enhanced_text.split("\n\n") if s.strip()]
    for section in sections:
        lines = section.split("\n")
        if not lines:
            continue
        
        slide_title = lines[0].replace("**Slide ", "").replace(":**", "").replace("**", "").replace("*", "").replace("\"","").strip()
        content = "\n".join(lines[1:]).strip()

        # Ignore the slide ONLY IF it contains redundant text AND nothing else
        if all(note in content for note in redundant_texts):
            continue

        add_slide(prs, slide_title, content)
        analyze_text_for_graphs(content, prs)

    if trend_data:
        create_trend_chart(prs, trend_data)

    # Process Tables for Charts
    table_index = 1
    if tables:
        for table in tables:
            # Check if the first row contains actual titles or just data
            if table.iloc[0].str.match(r'^[A-Za-z ]+$').all():  # Checks if all values in the first row are text
                title = table.iloc[0, 0]  # Use first column header as the title
                table = table.iloc[1:].reset_index(drop=True)  # Remove the first row if it's a title
            else:
                title = f"Table {table_index}"  # Assign a default title

            categories = list(table.iloc[:, 0])
            for col in table.columns[1:]:
                values = pd.to_numeric(table[col], errors="coerce").fillna(0).tolist()
                create_chart(prs, f"{title} - {col}", categories, values, XL_CHART_TYPE.COLUMN_CLUSTERED)
            table_index += 1


    # Process Excel Data - Only include if it has meaningful numerical content
    if excel_data is not None and isinstance(excel_data, pd.DataFrame):
        numeric_cols = excel_data.select_dtypes(include=["number"])
        text_cols = excel_data.select_dtypes(exclude=["number"])

        if not numeric_cols.empty or not text_cols.empty:  # Ensures slide is created if ANY relevant content exists
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Excel Data Analysis"
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.clear()

            # Add data content
            for index, row in excel_data.iterrows():
                p = text_frame.add_paragraph()
                p.text = " | ".join(map(str, row.values))
                p.font.size = SLIDE_FONT_SIZE

            # Generate Charts if numerical data is present
            if not numeric_cols.empty:
                categories = excel_data.iloc[:, 0].astype(str).tolist()
                for col in numeric_cols.columns:
                    values = excel_data[col].fillna(0).tolist()
                    create_chart(prs, f"{col} Analysis", categories, values, XL_CHART_TYPE.COLUMN_CLUSTERED)
    doc_title = re.sub(r'[<>:"/\\|?*]', '', title)  # Sanitize filename
    ppt_path = os.path.join(tempfile.gettempdir(), f"{doc_title}.pptx")
    ppt_path = os.path.join(tempfile.gettempdir(), doc_title)
    prs.save(ppt_path)
    return ppt_path

# Function to process uploaded files
def process_files(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        text, tables = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        text, tables = extract_text_from_docx(file_path)
    elif ext == ".pptx":
        text, tables = extract_text_from_pptx(file_path)
    elif ext in [".xls", ".xlsx"]:
        excel_data = extract_data_from_excel(file_path)
        return create_ppt_from_text("", None, excel_data)
    else:
        return None
    return create_ppt_from_text(text, tables)

# Streamlit UI
def main():
    plt.clf()
    st.title("Generate AI-Enhanced PPT from Documents")
    uploaded_file = st.file_uploader("Upload a Word, PPT, Excel, or PDF file", type=["pdf", "docx", "pptx", "xls", "xlsx"])
    
    if uploaded_file is not None:
        file_path = os.path.join(tempfile.gettempdir(), uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        ppt_path = process_files(file_path)
        
        with open(ppt_path, "rb") as f:
            st.download_button("Download AI-Enhanced PPT", f, file_name="Generated_Presentation.pptx")

if __name__ == "__main__":
    main() 